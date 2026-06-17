"""
slides_ppt.py — Converte o TXT de slides gerado pelo LLM em um arquivo PPTX.

Esperado no TXT:
  Slide 1 - Título:
  - bullet 1
  - bullet 2
  Slide 2 - Assunto:
  - bullet 1
  ...

Uso:
  slides_txt_to_ppt(
      slides_txt_path,          # Path para o .txt
      template_path=...,        # Path para o .pptx de template (opcional)
      images_dir=...,           # Path para pasta de imagens (opcional, não usado ainda)
  )
  → salva <slides_txt_path>.pptx ao lado do .txt
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# python-pptx é dependência obrigatória
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError as exc:  # pragma: no cover
    raise ImportError(
        "python-pptx não instalado. Execute: pip install python-pptx"
    ) from exc


# ---------------------------------------------------------------------------
# Parsing do TXT
# ---------------------------------------------------------------------------

_SLIDE_HEADER = re.compile(
    r"^Slide\s+(\d+)\s*[-–:]\s*(.+)$", re.IGNORECASE
)
_BULLET = re.compile(r"^[-*•]\s+(.+)$")


def parse_slides_txt(text: str) -> list[dict]:
    """
    Lê o texto gerado pelo LLM e retorna lista de dicts:
      [{"title": str, "bullets": [str, ...]}, ...]
    """
    slides: list[dict] = []
    current: Optional[dict] = None

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        header_match = _SLIDE_HEADER.match(line)
        if header_match:
            if current is not None:
                slides.append(current)
            title = header_match.group(2).strip().rstrip(":")
            current = {"title": title, "bullets": []}
            continue

        if current is not None:
            bullet_match = _BULLET.match(line)
            if bullet_match:
                current["bullets"].append(bullet_match.group(1).strip())
            elif line and not line.startswith("Slide"):
                # linha de texto sem marcador — trata como bullet simples
                current["bullets"].append(line)

    if current is not None:
        slides.append(current)

    return slides


# ---------------------------------------------------------------------------
# Montagem do PPTX
# ---------------------------------------------------------------------------

# Índices de layout padrão do PowerPoint (variam por template)
_LAYOUT_TITLE    = 0   # slide de título (capa)
_LAYOUT_CONTENT  = 1   # título + conteúdo


def _get_layout(prs: Presentation, index: int):
    """Retorna o slide layout pelo índice; usa o primeiro se o índice não existir."""
    layouts = prs.slide_layouts
    if index < len(layouts):
        return layouts[index]
    return layouts[0]


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    layout = _get_layout(prs, _LAYOUT_TITLE)
    slide = prs.slides.add_slide(layout)

    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    # idx 0 = título principal
    if 0 in placeholders:
        placeholders[0].text = title
    # idx 1 = subtítulo / corpo
    if 1 in placeholders and subtitle:
        placeholders[1].text = subtitle


def _add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
) -> None:
    layout = _get_layout(prs, _LAYOUT_CONTENT)
    slide = prs.slides.add_slide(layout)

    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    # Título
    if 0 in placeholders:
        placeholders[0].text = title

    # Corpo / bullets
    if 1 in placeholders and bullets:
        tf = placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = bullet
                tf.paragraphs[0].level = 0
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0


# ---------------------------------------------------------------------------
# Ponto de entrada público
# ---------------------------------------------------------------------------

def slides_txt_to_ppt(
    slides_txt_path: Path,
    *,
    template_path: Optional[Path] = None,
    images_dir: Optional[Path] = None,  # reservado para uso futuro
) -> Path:
    """
    Lê *slides_txt_path*, faz o parse e gera um PPTX ao lado com o mesmo stem.

    Parâmetros
    ----------
    slides_txt_path : Path
        Arquivo .txt com o conteúdo dos slides no formato padrão.
    template_path : Path, opcional
        Template .pptx a ser usado como base. Se None, cria uma apresentação em branco.
    images_dir : Path, opcional
        Pasta de imagens (reservado; ainda não usado).

    Retorna
    -------
    Path
        Caminho do arquivo .pptx gerado.
    """
    text = slides_txt_path.read_text(encoding="utf-8")
    slides_data = parse_slides_txt(text)

    if not slides_data:
        raise ValueError(
            f"Nenhum slide encontrado em {slides_txt_path}. "
            "Verifique se o formato é 'Slide N - Título:'"
        )

    # Carrega template ou cria apresentação vazia
    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
        # Remove slides pré-existentes do template, mantendo apenas master/layouts
        xml_slides = prs.slides._sldIdLst  # noqa: SLF001
        for _ in range(len(prs.slides)):
            rId = prs.slides._sldIdLst[0].get("r:id")  # noqa: SLF001
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]  # noqa: SLF001
        logger.info("Template carregado: %s", template_path.name)
    else:
        prs = Presentation()
        if template_path:
            logger.warning(
                "Template não encontrado (%s); usando apresentação em branco.",
                template_path,
            )

    # Adiciona os slides
    for i, slide in enumerate(slides_data):
        title   = slide["title"]
        bullets = slide["bullets"]

        if i == 0:
            # Primeiro slide: layout de capa
            subtitle = bullets[0] if bullets else ""
            _add_title_slide(prs, title, subtitle)
        else:
            _add_content_slide(prs, title, bullets)

    out_path = slides_txt_path.with_suffix(".pptx")
    prs.save(str(out_path))
    logger.info("PPTX salvo: %s (%d slides)", out_path.name, len(slides_data))
    return out_path
